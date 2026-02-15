from __future__ import annotations
import csv
import io
import os
import re
import tempfile
import subprocess
import uuid
from dataclasses import dataclass
from typing import Dict, List, Tuple, Optional

from django.conf import settings
from openpyxl import load_workbook
import qrcode
from pptx import Presentation

TOKEN_LIST = ["fname", "fschool", "fclass", "fplace", "fteacher", "fnomination", "fid"]

def mask_iin(iin: str) -> str:
    if not iin or len(iin) < 8:
        return ""
    return iin[:4] + "*" * 6 + iin[-2:]

def read_table(file_path: str) -> Tuple[List[str], List[Dict[str, str]]]:
    """Return (columns, rows) from CSV or XLSX."""
    lower = file_path.lower()
    if lower.endswith(".csv"):
        with open(file_path, "r", encoding="utf-8-sig", newline="") as f:
            reader = csv.DictReader(f)
            cols = reader.fieldnames or []
            rows = []
            for r in reader:
                rows.append({k: (v.strip() if isinstance(v, str) else v) for k, v in r.items()})
            return cols, rows

    if lower.endswith(".xlsx") or lower.endswith(".xlsm") or lower.endswith(".xltx") or lower.endswith(".xltm"):
        wb = load_workbook(file_path, read_only=True, data_only=True)
        ws = wb.active
        headers = []
        rows = []
        for i, row in enumerate(ws.iter_rows(values_only=True)):
            if i == 0:
                headers = [str(c).strip() if c is not None else "" for c in row]
                continue
            if all(c is None for c in row):
                continue
            data = {}
            for h, c in zip(headers, row):
                if not h:
                    continue
                data[h] = "" if c is None else str(c).strip()
            rows.append(data)
        return headers, rows

    raise ValueError("Поддерживаются только CSV и XLSX.")

def default_mapping(columns: List[str]) -> Dict[str, str]:
    """Auto map columns -> tokens."""
    cols = {c.strip().lower(): c for c in columns if c}
    mapping = {}
    # canonical
    # columns: name, id, iin, school, class, place, teacher, nomination
    def pick(*keys):
        for k in keys:
            if k in cols:
                return cols[k]
        return ""
    mapping["fname"] = pick("name", "fullname", "fio", "фио")
    mapping["fid"] = pick("id", "code", "номер", "номер/код диплома")
    mapping["fschool"] = pick("school", "school_name", "школа")
    mapping["fclass"] = pick("class", "grade", "класс")
    mapping["fplace"] = pick("place", "degree", "место", "степень")
    mapping["fteacher"] = pick("teacher", "teacher_name", "учитель")
    mapping["fnomination"] = pick("nomination", "category", "номинация")
    return mapping

def validate_required(columns: List[str]) -> List[str]:
    needed = ["name", "id", "iin", "school", "class", "place", "teacher", "nomination"]
    present = set(c.strip().lower() for c in columns)
    missing = [c for c in needed if c not in present]
    return missing

def _replace_in_paragraph_preserving_runs(paragraph, token_to_value: Dict[str, str]):
    """Replace tokens even if split across runs.
    Strategy:
      - Build full text + run spans
      - Find token occurrences in full text
      - For each occurrence: write replacement into first run, blank rest intersecting spans
    Formatting of first run is preserved.
    """
    runs = list(paragraph.runs)
    if not runs:
        return
    full = "".join(r.text for r in runs)
    if not full:
        return

    # Replace sequentially to keep spans manageable: do longest tokens first
    tokens = sorted(token_to_value.keys(), key=len, reverse=True)
    # We operate by scanning occurrences; for each occurrence, we mutate runs text and rebuild full/spans
    for tok in tokens:
        val = str(token_to_value.get(tok, ""))
        if not tok or tok not in full:
            continue

        # rebuild spans
        spans = []
        pos = 0
        for r in runs:
            spans.append((pos, pos + len(r.text)))
            pos += len(r.text)

        # find all occurrences of tok in current full
        start = 0
        while True:
            idx = full.find(tok, start)
            if idx == -1:
                break
            jdx = idx + len(tok)

            # locate first/last run indices that overlap [idx, jdx)
            first_i = None
            last_i = None
            for i, (a, b) in enumerate(spans):
                if b <= idx:
                    continue
                if a >= jdx:
                    break
                if first_i is None:
                    first_i = i
                last_i = i

            if first_i is None:
                break

            # Build new texts for affected runs
            # We'll keep prefix from first run up to idx, and suffix from last run after jdx
            prefix = full[:idx]
            suffix = full[jdx:]

            # Now reconstruct run texts with minimal change:
            # - Keep all runs before first_i as-is
            # - Put into first_i: its part of prefix within that run + replacement + (nothing else)
            # - Blank intermediate runs
            # - Put into last_i: keep suffix-part that belongs to last run? We'll do simpler:
            #     After replacement, we set all runs from first_i+1..last_i to "" and
            #     put the remainder (suffix) into last_i if last_i != first_i by appending to last_i.
            # But we must avoid moving too much text between runs; instead we rebuild whole paragraph text into first run
            # while preserving formatting: acceptable for MVP, but may lose mixed formatting inside one paragraph.
            # To reduce formatting loss, we keep runs before first_i and after last_i, and only rebuild within span.
            #
            # Extract left part within first run and right part within last run:
            first_a, first_b = spans[first_i]
            last_a, last_b = spans[last_i]

            left_in_first = full[first_a:idx]  # part of first run before token
            right_in_last = full[jdx:last_b]   # part of last run after token

            runs[first_i].text = left_in_first + val + right_in_last
            for k in range(first_i + 1, last_i + 1):
                if k == last_i:
                    # we've already kept right_in_last in first_i, so clear last_i
                    runs[k].text = ""
                else:
                    runs[k].text = ""

            # rebuild full and continue search after idx + len(val)
            full = "".join(r.text for r in runs)
            start = idx + len(val)

def replace_tokens_in_slide(prs: Presentation, token_to_value: Dict[str, str]):
    slide = prs.slides[0]
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        tf = shape.text_frame
        for p in tf.paragraphs:
            _replace_in_paragraph_preserving_runs(p, token_to_value)

def make_qr_png_bytes(url: str) -> bytes:
    img = qrcode.make(url)
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()

def insert_qr(prs: Presentation, qr_png_bytes: bytes) -> bool:
    """Try shape.name == 'QR' first, fallback to marker text 'fqr'.
    Returns True if inserted.
    """
    slide = prs.slides[0]
    # 1) shape name QR
    for shape in list(slide.shapes):
        if getattr(shape, "name", "") == "QR":
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            # remove placeholder
            el = shape._element
            el.getparent().remove(el)
            # insert picture
            with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                tmp.write(qr_png_bytes)
                tmp_path = tmp.name
            try:
                slide.shapes.add_picture(tmp_path, left, top, width=width, height=height)
            finally:
                try:
                    os.unlink(tmp_path)
                except OSError:
                    pass
            return True

    # 2) fallback: text marker fqr
    for shape in list(slide.shapes):
        if not getattr(shape, "has_text_frame", False):
            continue
        txt = shape.text_frame.text or ""
        if "fqr" in txt:
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            # remove marker shape (or we could clear text)
            el = shape._element
            el.getparent().remove(el)
            with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                tmp.write(qr_png_bytes)
                tmp_path = tmp.name
            try:
                slide.shapes.add_picture(tmp_path, left, top, width=width, height=height)
            finally:
                try:
                    os.unlink(tmp_path)
                except OSError:
                    pass
            return True
    return False

def convert_pptx_to_pdf(pptx_path: str, out_dir: str) -> str:
    """Convert via LibreOffice headless. Returns pdf path."""
    os.makedirs(out_dir, exist_ok=True)
    cmd = [
        "soffice",
        "--headless",
        "--nologo",
        "--nofirststartwizard",
        "--convert-to",
        "pdf",
        "--outdir",
        out_dir,
        pptx_path,
    ]
    proc = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
    if proc.returncode != 0:
        raise RuntimeError(f"LibreOffice failed: {proc.stderr or proc.stdout}")
    # LibreOffice writes pdf with same basename
    base = os.path.splitext(os.path.basename(pptx_path))[0]
    pdf_path = os.path.join(out_dir, base + ".pdf")
    if not os.path.exists(pdf_path):
        # sometimes produces upper/lower variants; search
        for f in os.listdir(out_dir):
            if f.lower().endswith(".pdf") and os.path.splitext(f)[0] == base:
                return os.path.join(out_dir, f)
        raise FileNotFoundError("PDF not found after conversion")
    return pdf_path
